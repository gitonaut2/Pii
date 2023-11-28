import os
from flask import Flask, request, render_template, send_file
from werkzeug.utils import secure_filename
import logging
from pptx import Presentation
from flask import send_file

logging.basicConfig(level=logging.INFO)
app = Flask(__name__)
app.secret_key = 'ein_zufälliger_schlüssel'
UPLOAD_FOLDER = r'C:\Users\admin\Desktop\App_OCR\app - Kopie\neue'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

BRAND_NAME_MAPPINGS = {
    'Makrosoft': 'Microsoft',
    'Lenevo': 'Lenovo',
    'Sorny': 'Sony',
    'Adibas': 'Adidas',
    'Sansung': 'Samsung'
}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'ppt', 'pptx'}

def analyze_pptx_and_find_corrections(pptx_path):
    prs = Presentation(pptx_path)
    corrections = {}
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                for fake_name, real_name in BRAND_NAME_MAPPINGS.items():
                    if fake_name in text:
                        corrections[fake_name] = real_name
    return corrections

def save_results_to_txt(corrections, filepath):
    with open(filepath, 'w') as file:
        for fake, real in corrections.items():
            file.write(f"{fake}: {real}\n")

@app.route('/', methods=['GET', 'POST'])
def upload_file():
    if request.method == 'POST':
        file = request.files['file']  # Achten Sie darauf, dass der Name 'file' im HTML-Formular übereinstimmt
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            saved_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(saved_path)

            if filename.endswith('.pptx'):
                corrections = analyze_pptx_and_find_corrections(saved_path)
                
                # Speichern der Ergebnisse in einer Textdatei
                txt_filename = 'Ergebnisse.txt'
                temp_txt_path = os.path.join(app.config['UPLOAD_FOLDER'], txt_filename)
                save_results_to_txt(corrections, temp_txt_path)

                # Aufräumen: Originaldatei entfernen
                os.remove(saved_path)

                # Schicken Sie sowohl die Korrekturen als auch den Dateinamen der Ergebnisdatei an das Template
                return render_template('results.html', message=corrections, filename=txt_filename)
            else:
                os.remove(saved_path)  # Aufräumen: Ungültige Datei entfernen
                return 'Ungültiges Dateiformat', 400
            
    # Zeigen Sie das Upload-Formular an, wenn keine POST-Anfrage vorliegt
    return render_template('upload.html')


@app.route('/download_results', methods=['POST'])
def download_results():
    filename = request.form['filename']
    if filename:
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        try:
            return send_file(filepath, as_attachment=True, attachment_filename=filename)
        except Exception as e:
            logging.error(f"An error occurred: {e}")
            return str(e), 500
    else:
        return 'Fehler: Kein Dateiname angegeben', 400


if __name__ == "__main__":
    app.run(debug=True)
